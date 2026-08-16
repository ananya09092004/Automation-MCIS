"""Reliable Office-file capability pack without fragile Windows COM automation.

Files are created and verified with their native document formats.  Opening a
file still uses the user's installed default Office app through ``os.startfile``.
"""

import os
from pathlib import Path


class OfficeUnavailableError(RuntimeError):
    pass


class OfficeCapability:
    @staticmethod
    def _output(path: str) -> Path:
        output = Path(path).resolve()
        output.parent.mkdir(parents=True, exist_ok=True)
        return output

    @staticmethod
    def _dependency(package: str, error: Exception):
        raise OfficeUnavailableError(f"Install {package} to use this Office file capability.") from error

    def create_word_document(self, path: str, text: str = "") -> str:
        try:
            from docx import Document
        except ImportError as error:
            self._dependency("python-docx", error)
        output = self._output(path)
        document = Document()
        if text:
            document.add_paragraph(text)
        document.save(output)
        if not output.is_file():
            raise RuntimeError("Word document was not created.")
        return str(output)

    def read_word_document(self, path: str) -> str:
        try:
            from docx import Document
        except ImportError as error:
            self._dependency("python-docx", error)
        document = Document(str(Path(path).resolve()))
        return "\n".join(paragraph.text for paragraph in document.paragraphs).strip()

    def create_excel_workbook(self, path: str, rows: list[list] | None = None) -> str:
        try:
            from openpyxl import Workbook
        except ImportError as error:
            self._dependency("openpyxl", error)
        output = self._output(path)
        workbook = Workbook()
        sheet = workbook.active
        for row in rows or []:
            sheet.append(list(row))
        workbook.save(output)
        workbook.close()
        if not output.is_file():
            raise RuntimeError("Excel workbook was not created.")
        return str(output)

    def read_excel_rows(self, path: str, max_rows: int = 100, max_columns: int = 50) -> list[list]:
        try:
            from openpyxl import load_workbook
        except ImportError as error:
            self._dependency("openpyxl", error)
        workbook = load_workbook(Path(path).resolve(), read_only=True, data_only=True)
        try:
            sheet = workbook.active
            return [list(row[:max_columns]) for row in sheet.iter_rows(max_row=max_rows, values_only=True)
                    if any(value is not None for value in row)]
        finally:
            workbook.close()

    def create_powerpoint_presentation(self, path: str, title: str = "", body: str = "") -> str:
        try:
            from pptx import Presentation
        except ImportError as error:
            self._dependency("python-pptx", error)
        output = self._output(path)
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[0])
        if title:
            slide.shapes.title.text = title
        if body and len(slide.placeholders) > 1:
            slide.placeholders[1].text = body
        presentation.save(output)
        if not output.is_file():
            raise RuntimeError("PowerPoint presentation was not created.")
        return str(output)

    def inspect_powerpoint_presentation(self, path: str) -> dict:
        try:
            from pptx import Presentation
        except ImportError as error:
            self._dependency("python-pptx", error)
        presentation = Presentation(Path(path).resolve())
        return {"slide_count": len(presentation.slides)}

    def open_document(self, path: str) -> bool:
        target = Path(path).resolve()
        if not target.is_file():
            return False
        os.startfile(str(target))
        return True
